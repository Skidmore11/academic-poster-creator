#!/usr/bin/env python3
"""
Academic Poster Creator - Web Interface
A modern web app to upload PDFs, extract information using AI, and create academic posters.
"""

import os
import tempfile
import shutil
from flask import Flask, render_template, request, redirect, url_for, flash, send_file, jsonify
from werkzeug.utils import secure_filename
import pdfplumber  # PDF text extraction
import openai
import anthropic
from pptx import Presentation
from pptx.util import Pt, Inches
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
import re
import json
from datetime import datetime
import random
from dotenv import load_dotenv
from pptx.enum.shapes import MSO_SHAPE_TYPE
import template_configs

# Load environment variables from .env file
load_dotenv()

# ============================================================================
# CONFIGURATION - EDIT THESE SETTINGS
# ============================================================================

# 🔑 API Configuration
OPENAI_API_KEY = os.getenv('OPENAI_API_KEY')
ANTHROPIC_API_KEY = os.getenv('ANTHROPIC_API_KEY')

# Default API provider (can be changed via frontend)
DEFAULT_API_PROVIDER = 'openai'  # 'openai' or 'anthropic'

# Global variable to track current API provider
current_api_provider = DEFAULT_API_PROVIDER

# Validate API keys (at least one must be set)
if not OPENAI_API_KEY and not ANTHROPIC_API_KEY:
    raise ValueError("At least one API key must be set. Please set OPENAI_API_KEY or ANTHROPIC_API_KEY in your .env file.")

# 📁 File Settings
UPLOAD_FOLDER = 'uploads'
TEMPLATE_LIBRARY_FOLDER = 'template_library'
ALLOWED_EXTENSIONS = {'pdf', 'png', 'jpg', 'jpeg', 'pptx'}
MAX_CONTENT_LENGTH = 500 * 1024 * 1024  # 500MB max file size

# 🧹 Cleanup Settings - Set to True to automatically delete uploaded files after processing
AUTO_CLEANUP_UPLOADS = True
KEEP_FINAL_OUTPUT = True  # Keep the final PowerPoint file for download

# 🎨 Web App Settings
app = Flask(__name__)
app.secret_key = os.getenv('FLASK_SECRET_KEY', os.urandom(24))
app.config['MAX_CONTENT_LENGTH'] = MAX_CONTENT_LENGTH
app.config['MAX_CONTENT_PATH'] = None

# 🧪 Testing Settings - Set to True to use dummy data instead of API calls
USE_DUMMY_DATA = False
DUMMY_DATA_FILE = 'dummy_api_response.json'

# Global variable to track current mode (can be changed via API)
current_dummy_mode = USE_DUMMY_DATA

# ============================================================================
# END CONFIGURATION - DON'T EDIT BELOW THIS LINE
# ============================================================================

# Create upload folder if it doesn't exist
os.makedirs(UPLOAD_FOLDER, exist_ok=True)
os.makedirs(TEMPLATE_LIBRARY_FOLDER, exist_ok=True)

def cleanup_uploaded_files(files_to_cleanup, keep_final_output=True):
    """
    Clean up uploaded files after processing.
    
    Args:
        files_to_cleanup (list): List of file paths to delete
        keep_final_output (bool): Whether to keep the final PowerPoint output file
    """
    if not AUTO_CLEANUP_UPLOADS:
        return
    
    for file_path in files_to_cleanup:
        try:
            # Skip the final output file if keep_final_output is True
            if keep_final_output and file_path.endswith('.pptx') and '_academic_' in file_path:
                print(f"🔄 Keeping final output: {os.path.basename(file_path)}")
                continue
                
            if os.path.exists(file_path):
                os.remove(file_path)
                print(f"🗑️ Cleaned up: {os.path.basename(file_path)}")
        except Exception as e:
            print(f"⚠️ Warning: Could not delete {file_path}: {e}")

def cleanup_old_files(days_old=1):
    """
    Clean up files older than specified days in uploads folder.
    Only runs if AUTO_CLEANUP_UPLOADS is True.
    """
    if not AUTO_CLEANUP_UPLOADS:
        return
        
    try:
        current_time = datetime.now()
        cutoff_time = current_time.replace(hour=0, minute=0, second=0, microsecond=0)
        
        for filename in os.listdir(UPLOAD_FOLDER):
            file_path = os.path.join(UPLOAD_FOLDER, filename)
            if os.path.isfile(file_path):
                file_time = datetime.fromtimestamp(os.path.getmtime(file_path))
                if file_time < cutoff_time:
                    try:
                        os.remove(file_path)
                        print(f"🗑️ Cleaned up old file: {filename}")
                    except Exception as e:
                        print(f"⚠️ Could not delete old file {filename}: {e}")
    except Exception as e:
        print(f"⚠️ Error during cleanup: {e}")

def save_dummy_data(extracted_data):
    """Save API response as dummy data for future testing."""
    try:
        dummy_data = {
            'timestamp': datetime.now().isoformat(),
            'data': extracted_data
        }
        with open(DUMMY_DATA_FILE, 'w', encoding='utf-8') as f:
            json.dump(dummy_data, f, indent=2, ensure_ascii=False)
        print(f"✅ Dummy data saved to {DUMMY_DATA_FILE}")
        return True
    except Exception as e:
        print(f"❌ Error saving dummy data: {e}")
        return False

def load_dummy_data():
    """Load dummy data for testing without API calls."""
    try:
        if os.path.exists(DUMMY_DATA_FILE):
            with open(DUMMY_DATA_FILE, 'r', encoding='utf-8') as f:
                dummy_data = json.load(f)
            print(f"✅ Loaded dummy data from {DUMMY_DATA_FILE}")
            return dummy_data['data'], None
        else:
            return None, "No dummy data file found"
    except Exception as e:
        return None, f"Error loading dummy data: {e}"

def load_template_library():
    """Get list of templates in the library using folder-based status detection."""
    try:
        templates = []
        
        # Define folder structure for template status
        folder_structure = {
            'available': os.path.join(TEMPLATE_LIBRARY_FOLDER, 'available'),
            'coming_soon': os.path.join(TEMPLATE_LIBRARY_FOLDER, 'coming_soon'),
            'premium': os.path.join(TEMPLATE_LIBRARY_FOLDER, 'premium')
        }
        
        # Process each folder
        for status, folder_path in folder_structure.items():
            if os.path.exists(folder_path):
                for filename in os.listdir(folder_path):
                    if filename.endswith('.pptx'):
                        template_path = os.path.join(folder_path, filename)
                        file_size = os.path.getsize(template_path)
                        file_size_mb = round(file_size / (1024 * 1024), 2)
                        
                        # Prefer manual preview if available
                        base = os.path.splitext(filename)[0]
                        manual_preview = None
                        for ext in ['.png', '.jpg', '.jpeg']:
                            candidate = base + '_manual_preview' + ext
                            candidate_path = os.path.join(folder_path, candidate)
                            if os.path.exists(candidate_path):
                                manual_preview = candidate
                                break
                        
                        # Otherwise, generate or use auto preview
                        if manual_preview:
                            preview_filename = manual_preview
                        else:
                            preview_filename = base + '_preview.png'
                            preview_path = os.path.join(folder_path, preview_filename)
                            if not os.path.exists(preview_path):
                                generate_template_preview(template_path, preview_path)
                            if not os.path.exists(preview_path):
                                preview_filename = None
                        
                        # Get template name and determine status based on folder
                        template_name = base.replace('_', ' ').title()
                        
                        # Determine status based on folder location
                        is_premium = (status == 'premium')
                        is_coming_soon = (status == 'coming_soon')
                        is_new = False  # Could be determined by file creation date or a "new" folder
                        
                        templates.append({
                            'filename': filename,
                            'name': template_name,
                            'size_mb': file_size_mb,
                            'path': template_path,
                            'preview': preview_filename,
                            'is_premium': is_premium,
                            'is_coming_soon': is_coming_soon,
                            'is_new': is_new,
                            'folder': status  # Add folder info for debugging
                        })
        
        # Sort by name
        templates.sort(key=lambda x: x['name'])
        return templates
    except Exception as e:
        print(f"❌ Error loading template library: {e}")
        return []

def save_template_to_library(template_file, filename, folder='available'):
    """Save uploaded template to the library in the specified folder."""
    try:
        # Create the target folder if it doesn't exist
        target_folder = os.path.join(TEMPLATE_LIBRARY_FOLDER, folder)
        os.makedirs(target_folder, exist_ok=True)
        
        template_path = os.path.join(target_folder, filename)
        template_file.save(template_path)
        print(f"✅ Template saved to library ({folder}): {filename}")
        return True, None
    except Exception as e:
        return False, f"Error saving template: {e}"

def generate_template_preview(template_path, preview_path):
    """Generate a preview image from PowerPoint template."""
    try:
        from pptx import Presentation
        
        # Load the presentation
        prs = Presentation(template_path)
        if len(prs.slides) == 0:
            return False, "No slides found in template"
        
        # For now, just return success without generating preview
        # This avoids the Pillow dependency
        return True, None
        
    except Exception as e:
        return False, f"Error generating preview: {e}"

def allowed_file(filename, extensions):
    """Check if file extension is allowed."""
    return '.' in filename and filename.rsplit('.', 1)[1].lower() in extensions

def extract_text_from_pdf(file_path):
    """Extract text from PDF file using pdfplumber."""
    try:
        with pdfplumber.open(file_path) as pdf:
            text = ""
            for page in pdf.pages:
                page_text = page.extract_text()
                if page_text:
                    text += page_text + "\n"
            return text
    except Exception as e:
        return f"Error extracting text from PDF: {e}"

def call_ai_api(manuscript_text, provider='openai'):
    """Call the specified AI API to extract poster content."""
    global current_api_provider
    
    # Use the specified provider or fall back to current global setting
    if not provider:
        provider = current_api_provider
    
    print(f"🤖 Using {provider.upper()} API for content extraction...")
    
    # Common prompt for all APIs
    prompt = f"""
You are an expert in academic writing and research poster design.

Given the following research manuscript text, extract content for an academic A0-size research poster. 

**Instructions:**
- Return your answer as a single, valid JSON object (not Python, not markdown, not prose).
- Use only double quotes for all keys and values.
- Do not include any introductory or closing text, explanations, or markdown formatting.
- Each value should be a single line (no line breaks inside values).
- If a value is missing, use an empty string.
- Do not include any keys other than those listed below.
- **CRITICAL: Include COMPLETE author lists and affiliations - DO NOT use "et al" or "(see manuscript for full list)".**
- **CRITICAL: Count words carefully and stick to the exact word count ranges specified for each field.**
- **Word count includes all words, including articles (a, an, the) and prepositions.**
- **If a field is too short, expand it with more detail. If too long, condense it while keeping key information.**

**JSON keys and requirements:**
- "headline": A short, punchy phrase (3-8 words) summarizing the main finding. Surround 2-5 important words with asterisks (e.g., *BOOSTS* or *DIGITAL HEALTH*).
- "title": The full title of the research.
- "authors": **COMPLETE list of all authors** - include every author name found in the manuscript. Do not use "et al" or truncate the list.
- "affiliations": **COMPLETE list of all affiliations** - include every institution, department, and affiliation mentioned. Do not use "(see manuscript for full list)" or truncate.
- "subtitle": (optional) A brief subtitle if available.
- "Introduction": EXACTLY 65-75 words. Provide a comprehensive background, context, and rationale for the study. Include key concepts, current state of knowledge, and gaps that justify the research. PRESERVE THE EXACT CITATION STYLE from the original manuscript (e.g., if the PDF uses "(Smith et al., 2020)" or "[1]" or "¹", keep that exact format).
- "Objective": 15-25 words. Clear, specific research objective or question.
- "Methods": EXACTLY 80-90 words. Include study design, participants, procedures, data collection, and analysis methods. Be specific about sample size, timeframes, and key variables. PRESERVE THE EXACT CITATION STYLE from the original manuscript.
- "Results": EXACTLY 80-90 words. Present key findings with specific numbers, percentages, or statistics when available. Include sample sizes, effect sizes, and significance levels.
- "Discussion": 60-80 words. Interpret the main findings, discuss implications, limitations, and future directions. PRESERVE THE EXACT CITATION STYLE from the original manuscript.
- "Conclusions": EXACTLY 65-75 words. Summarize key findings and their significance. Include clinical or practical implications and recommendations.
- "References": Only references actually cited in the Introduction, Methods, or Discussion, PRESERVING THE EXACT REFERENCE FORMAT from the original manuscript (e.g., if the PDF uses APA style, keep APA style; if it uses Vancouver style, keep Vancouver style). Include complete reference details - do not use "et al" in references.

**Example output:**
{{
  "headline": "DIGITAL HEALTH *BOOSTS* OUTCOMES in chronic pain",
  "title": "Digital health interventions for chronic pain: A systematic review",
  "authors": "Smith J, Doe A, Johnson B, Williams C, Brown D, Davis E, Wilson F, Anderson G, Taylor H, Martinez I",
  "affiliations": "Department of Pain Medicine, University of Example; School of Health Sciences, Medical College; Institute of Digital Health, Technology University; Department of Psychology, State University; Center for Chronic Pain Research, National Institute",
  "subtitle": "",
  "Introduction": "Chronic pain affects millions of people worldwide (Cohen et al., 2021) and remains a significant public health challenge with substantial economic and social costs. Current treatment approaches often provide limited relief, creating an urgent need for innovative solutions. Digital health interventions, including mobile applications, wearable devices, and telehealth platforms, have emerged as promising alternatives that can deliver personalized care remotely (Fishman, 2021). This systematic review examines the effectiveness of these digital interventions in managing chronic pain conditions.",
  "Objective": "To evaluate the effectiveness of digital health interventions for chronic pain management.",
  "Methods": "We conducted a systematic review of randomized controlled trials published between 2010 and 2023. Electronic databases including PubMed, Embase, Cochrane Library, and PsycINFO were searched using relevant keywords. Studies were included if they evaluated digital interventions for chronic pain in adults. Primary outcomes were pain intensity and quality of life measures. Two independent reviewers screened articles and extracted data.",
  "Results": "Twenty-three studies met inclusion criteria with a total of 2,847 participants. Digital interventions led to significant reductions in pain intensity compared to usual care (mean difference -1.2 points on 0-10 scale, 95% CI -1.8 to -0.6). Quality of life improvements were also observed across multiple domains. Mobile applications showed the strongest effects, with 65% of studies reporting clinically meaningful improvements.",
  "Discussion": "Digital health interventions demonstrate promising results for chronic pain management, particularly mobile applications and telehealth platforms (Smith et al., 2022). However, heterogeneity in intervention types and outcome measures limits generalizability. Long-term effectiveness and cost-effectiveness require further investigation.",
  "Conclusions": "Digital health interventions can significantly improve outcomes for chronic pain patients, with mobile applications showing particular promise. These findings support the integration of digital solutions into pain management protocols. Future research should focus on long-term efficacy, cost-effectiveness, and implementation strategies.",
  "References": "Cohen SP, Vase L, Hooten WM. Chronic pain: an update on burden, best practices, and new advances. Lancet. 2021;397(10289):2082-2097. Fishman SM. Addressing the opioid crisis through education. Pain Med. 2021;22(4):741-742. Smith J, Doe A, Johnson B. Digital health interventions for chronic pain. Pain. 2022;163(5):1001-1010."
}}

**TEXT:**
{manuscript_text[:30000]}
"""

    try:
        if provider == 'openai':
            if not OPENAI_API_KEY:
                return None, "OpenAI API key not configured"
            
            print(f"🔧 Creating OpenAI client...")
            try:
                # Create a completely clean environment for client creation
                import os
                import copy
                import subprocess
                import sys
                
                # Save all current environment variables
                original_env = copy.deepcopy(os.environ)
                
                # Clear ALL environment variables that might interfere
                problematic_vars = [
                    'HTTP_PROXY', 'HTTPS_PROXY', 'http_proxy', 'https_proxy',
                    'NO_PROXY', 'no_proxy', 'ALL_PROXY', 'all_proxy',
                    'REQUESTS_CA_BUNDLE', 'CURL_CA_BUNDLE', 'SSL_CERT_FILE',
                    'REQUESTS_PROXIES', 'REQUESTS_VERIFY', 'REQUESTS_CERT',
                    'CURL_PROXY', 'CURL_VERIFY', 'CURL_CERT'
                ]
                
                cleared_vars = []
                for var in problematic_vars:
                    if var in os.environ:
                        cleared_vars.append(var)
                        del os.environ[var]
                        print(f"🔧 Cleared environment variable: {var}")
                
                # Also clear any requests-related environment variables
                for key in list(os.environ.keys()):
                    if 'proxy' in key.lower() or 'cert' in key.lower() or 'verify' in key.lower():
                        if key not in cleared_vars:
                            cleared_vars.append(key)
                            del os.environ[key]
                            print(f"🔧 Cleared additional environment variable: {key}")
                
                try:
                    # Create client with minimal environment
                    openai.api_key = OPENAI_API_KEY
                    print(f"✅ OpenAI client created successfully")
                finally:
                    # Restore original environment
                    for var in cleared_vars:
                        if var in original_env:
                            os.environ[var] = original_env[var]
                            print(f"🔧 Restored environment variable: {var}")
                        
            except Exception as e:
                print(f"❌ Error creating OpenAI client: {e}")
                print(f"🔧 Full error details: {type(e).__name__}: {str(e)}")
                print(f"🔧 Error type: {type(e)}")
                print(f"🔧 Error args: {e.args}")
                return None, f"Error creating OpenAI client: {e}"
            response = openai.ChatCompletion.create(
                model="gpt-4",
                messages=[
                    {"role": "system", "content": "You are an expert in academic writing and research poster design."},
                    {"role": "user", "content": prompt}
                ],
                temperature=0.4
            )
            raw_content = response.choices[0].message.content
            
        elif provider == 'anthropic':
            if not ANTHROPIC_API_KEY:
                return None, "Anthropic API key not configured"
            
            print(f"🔧 Creating Anthropic client...")
            try:
                # Create a completely clean environment for client creation
                import os
                import copy
                import subprocess
                import sys
                
                # Save all current environment variables
                original_env = copy.deepcopy(os.environ)
                
                # Clear ALL environment variables that might interfere
                problematic_vars = [
                    'HTTP_PROXY', 'HTTPS_PROXY', 'http_proxy', 'https_proxy',
                    'NO_PROXY', 'no_proxy', 'ALL_PROXY', 'all_proxy',
                    'REQUESTS_CA_BUNDLE', 'CURL_CA_BUNDLE', 'SSL_CERT_FILE',
                    'REQUESTS_PROXIES', 'REQUESTS_VERIFY', 'REQUESTS_CERT',
                    'CURL_PROXY', 'CURL_VERIFY', 'CURL_CERT'
                ]
                
                cleared_vars = []
                for var in problematic_vars:
                    if var in os.environ:
                        cleared_vars.append(var)
                        del os.environ[var]
                        print(f"🔧 Cleared environment variable: {var}")
                
                # Also clear any requests-related environment variables
                for key in list(os.environ.keys()):
                    if 'proxy' in key.lower() or 'cert' in key.lower() or 'verify' in key.lower():
                        if key not in cleared_vars:
                            cleared_vars.append(key)
                            del os.environ[key]
                            print(f"🔧 Cleared additional environment variable: {key}")
                
                try:
                    # Create client with minimal environment
                    anthropic.api_key = ANTHROPIC_API_KEY
                    print(f"✅ Anthropic client created successfully")
                finally:
                    # Restore original environment
                    for var in cleared_vars:
                        if var in original_env:
                            os.environ[var] = original_env[var]
                            print(f"🔧 Restored environment variable: {var}")
                        
            except Exception as e:
                print(f"❌ Error creating Anthropic client: {e}")
                print(f"🔧 Full error details: {type(e).__name__}: {str(e)}")
                print(f"🔧 Error type: {type(e)}")
                print(f"🔧 Error args: {e.args}")
                return None, f"Error creating Anthropic client: {e}"
            response = anthropic.Client().messages.create(
                model="claude-3-sonnet-20240229",
                max_tokens=4000,
                temperature=0.4,
                system="You are an expert in academic writing and research poster design.",
                messages=[
                    {"role": "user", "content": prompt}
                ]
            )
            raw_content = response.content[0].text
            
        else:
            return None, f"Unsupported API provider: {provider}"
        
        # Clean the response
        cleaned_content = (
            raw_content.strip()
            .removeprefix("```python")
            .removesuffix("```")
            .removeprefix("```json")
            .removesuffix("```")
            .strip()
        )
        
        print(f"🔍 Raw content length: {len(raw_content)}")
        print(f"🔍 Cleaned content length: {len(cleaned_content)}")
        print(f"🔍 First 200 chars: {cleaned_content[:200]}")
        
        # Try to extract JSON from text if it's wrapped
        if (cleaned_content.startswith("Here's the structured content") or 
            cleaned_content.startswith("Here's the formatted content") or
            cleaned_content.startswith("Here's the extracted content")):
            # Extract the JSON part after the colon
            json_start = cleaned_content.find('{')
            if json_start != -1:
                cleaned_content = cleaned_content[json_start:]
                print(f"🔍 Extracted JSON from text wrapper")
        
        # Clean up invalid control characters and line breaks in JSON strings
        import re
        # Replace newlines and carriage returns within quoted strings with spaces
        cleaned_content = re.sub(r'"([^"]*?)(?:\n|\r)([^"]*?)"', r'"\1 \2"', cleaned_content)
        # Also handle multi-line strings more comprehensively
        cleaned_content = re.sub(r'([^"])\n([^"])', r'\1 \2', cleaned_content)
        
        print(f"🔍 Final cleaned content length: {len(cleaned_content)}")
        print(f"🔍 Last 200 chars: {cleaned_content[-200:]}")
        
        # Parse the response
        try:
            # First try to parse as JSON
            try:
                import json
                poster = json.loads(cleaned_content)
                print(f"✅ Successfully parsed as JSON")
            except json.JSONDecodeError as json_error:
                print(f"⚠️ JSON parsing failed: {json_error}")
                # If JSON fails, try Python eval (handles single quotes)
                try:
                    poster = eval(cleaned_content)
                    print(f"✅ Successfully parsed with eval")
                except Exception as eval_error:
                    print(f"❌ Eval parsing failed: {eval_error}")
                    raise eval_error
            
            # Clean up references
            if 'References' in poster:
                ref = poster['References'].strip()
                print(f"🔍 [DEBUG] Raw references from API: '{ref}'")
                print(f"🔍 [DEBUG] References length: {len(ref)} characters")
                cleaned = ref.replace('[Reference details not found]', '').replace(' ,', ',').replace(', ,', ',').strip(',; .\n')
                print(f"🔍 [DEBUG] Cleaned references: '{cleaned}'")
                if cleaned and cleaned.strip():
                    poster['References'] = cleaned.strip()
                    print(f"✅ [DEBUG] Final references set: '{poster['References']}'")
                else:
                    poster['References'] = '[Reference details not found]'
                    print(f"⚠️ [DEBUG] References were empty, set to placeholder")
            else:
                print(f"⚠️ [DEBUG] No 'References' key found in API response")
                print(f"🔍 [DEBUG] Available keys: {list(poster.keys())}")
            
            return poster, None
            
        except Exception as parse_error:
            print(f"❌ Final parsing error: {parse_error}")
            return None, f"{provider.upper()} API response was not valid Python/JSON. Response was: {cleaned_content}"
            
    except Exception as e:
        return None, f"Error calling {provider.upper()} API: {e}"

def extract_information_from_pdf(manuscript_text):
    """Generate poster content using AI API or dummy data."""
    global current_dummy_mode, current_api_provider
    
    # Check if we should use dummy data
    if current_dummy_mode:
        print("🧪 Using dummy data instead of API call...")
        dummy_data, error = load_dummy_data()
        if error:
            return None, f"Dummy data error: {error}"
        return dummy_data, None
    
    # Call the AI API with current provider
    poster, error = call_ai_api(manuscript_text, current_api_provider)
    if error:
        return None, error
    
    # Save the API response as dummy data for future testing
    save_dummy_data(poster)
    return poster, None

def extract_information_from_pdf_with_provider(manuscript_text, provider):
    """Generate poster content using AI API with specified provider or dummy data."""
    global current_dummy_mode
    
    # Check if we should use dummy data
    if current_dummy_mode:
        print("🧪 Using dummy data instead of API call...")
        dummy_data, error = load_dummy_data()
        if error:
            return None, f"Dummy data error: {error}"
        return dummy_data, None
    
    # Call the AI API with specified provider
    poster, error = call_ai_api(manuscript_text, provider)
    if error:
        return None, error
    
    # Save the API response as dummy data for future testing
    save_dummy_data(poster)
    return poster, None

def find_shape_in_groups(slide, target_name):
    """Find shape by name including in groups."""
    for shape in slide.shapes:
        if shape.name.lower() == target_name.lower():
            return shape
        if hasattr(shape, 'shapes'):
            for subshape in shape.shapes:
                if subshape.name.lower() == target_name.lower():
                    return subshape
    return None

def get_title_font_size(title, template_name=None):
    """Return font size (Pt) for title based on content length and template configuration."""
    from template_configs import get_dynamic_font_size_config, get_default_dynamic_font_sizes
    
    length = len(title)
    
    # Get template-specific configuration
    if template_name:
        config = get_dynamic_font_size_config(template_name, "title")
    else:
        config = get_default_dynamic_font_sizes()["title"]
    
    # Apply template-specific sizing
    if length > 120:
        return Pt(config["long"])
    elif length > 80:
        return Pt(config["medium"])
    else:
        return Pt(config["short"])

def get_subtitle_font_size(title_font_size_pt, template_name=None):
    """Return subtitle font size (Pt) based on title size and template configuration."""
    from template_configs import get_dynamic_font_size_config, get_default_dynamic_font_sizes
    
    # Get template-specific configuration
    if template_name:
        config = get_dynamic_font_size_config(template_name, "subtitle")
    else:
        config = get_default_dynamic_font_sizes()["subtitle"]
    
    # Apply template-specific ratio and minimum size
    subtitle_size = int(round(title_font_size_pt * config["ratio"]))
    return Pt(max(subtitle_size, config["min_size"]))

def get_dynamic_font_size(text, template_name=None, section_type="body"):
    """Return font size (Pt) for text boxes based on text length and template configuration."""
    from template_configs import get_dynamic_font_size_config, get_default_dynamic_font_sizes
    
    length = len(text)
    
    # Get template-specific configuration
    if template_name:
        config = get_dynamic_font_size_config(template_name, section_type)
    else:
        default_config = get_default_dynamic_font_sizes()
        config = default_config.get("main_body", default_config["authors"])  # Fallback
    
    # Apply template-specific sizing
    if length <= 80:
        return Pt(config["short"])
    elif length <= 160:
        return Pt(config["medium"])
    elif length <= 250:
        return Pt(config["long"])
    else:
        # Check if extra_long is available, otherwise fall back to long
        if "extra_long" in config:
            return Pt(config["extra_long"])
        else:
            return Pt(config["long"])

def get_fixed_font_size(template_name=None, section_type="main_body_text"):
    """Return fixed font size (Pt) for sections that should not change based on content length."""
    from template_configs import get_dynamic_font_size_config, get_default_dynamic_font_sizes
    
    # Get template-specific configuration
    if template_name:
        config = get_dynamic_font_size_config(template_name, section_type)
    else:
        default_config = get_default_dynamic_font_sizes()
        config = default_config.get("main_body", default_config["authors"])  # Fallback
    
    # Use the "short" size as the standard fixed size
    return Pt(config["short"])

def calculate_image_fit(image_path, placeholder_width, placeholder_height):
    """Calculate the best fit for an image within placeholder dimensions while maintaining aspect ratio."""
    try:
        # For now, use placeholder dimensions directly
        # This avoids the Pillow dependency for image dimension calculation
        print(f"[DEBUG] Using placeholder dimensions: {placeholder_width}x{placeholder_height}")
        
        # Use placeholder dimensions with no offset (centered)
        return placeholder_width, placeholder_height, 0, 0
        
    except Exception as e:
        print(f"[DEBUG] Error calculating image fit: {e}")
        # Fallback to placeholder dimensions
        return placeholder_width, placeholder_height, 0, 0

def validate_image_file(image_path):
    """Validate that the image file is readable and has reasonable dimensions."""
    try:
        # Simple file size check instead of image dimensions
        file_size = os.path.getsize(image_path)
        
        # Check if file is too small (less than 1KB)
        if file_size < 1024:
            print(f"[WARNING] Image file is very small: {file_size} bytes")
            return False, f"Image file is too small ({file_size} bytes). Minimum size is 1KB."
        
        # Check if file is too large (more than 50MB)
        if file_size > 50 * 1024 * 1024:
            print(f"[WARNING] Image file is very large: {file_size} bytes")
            return False, f"Image file is too large ({file_size // (1024*1024)}MB). Maximum size is 50MB."
        
        print(f"[DEBUG] Image validation passed: {file_size} bytes")
        return True, None
        
    except Exception as e:
        return False, f"Error validating image: {e}"

def insert_image_safely(slide, image_path, placeholder_shape, placeholder_name):
    """Safely insert an image into a slide, always removing the placeholder and inserting the image as a new shape with calculated fit."""
    try:
        print(f"[DEBUG] {placeholder_name}: Inserting image with strict fit.")
        # Store placeholder dimensions and position
        left = placeholder_shape.left
        top = placeholder_shape.top
        width = placeholder_shape.width
        height = placeholder_shape.height

        # Calculate best fit for the image
        new_width, new_height, offset_x, offset_y = calculate_image_fit(image_path, width, height)

        # Adjust position to center the image
        adjusted_left = left + offset_x
        adjusted_top = top + offset_y

        # Remove the placeholder shape (always)
        try:
            sp = placeholder_shape._element
            sp.getparent().remove(sp)
        except Exception as e:
            print(f"[DEBUG] Could not remove placeholder shape: {e}")

        # Insert the image with calculated dimensions
        try:
            slide.shapes.add_picture(image_path, adjusted_left, adjusted_top, new_width, new_height)
            print(f"[DEBUG] Image inserted into {placeholder_name} with strict fit.")
            return True, None
        except Exception as e:
            print(f"[DEBUG] Error inserting image: {e}")
            return False, f"Error inserting image: {e}"

    except Exception as e:
        return False, f"Error in image insertion process: {e}"

def insert_colored_headline(shape, headline, highlight_rgb=(255,140,0), default_rgb=(255,255,255), font_family="Intro Rust", font_size=100, font_bold=True, alignment=None):
    # Remove all text first
    shape.text = ""
    # Split headline into normal and highlighted parts
    parts = re.split(r'(\*[^*]+\*)', headline)
    p = shape.text_frame.paragraphs[0]
    p.clear()
    if alignment is not None:
        p.alignment = alignment
    for part in parts:
        if not part:
            continue
        run = p.add_run()
        if part.startswith("*") and part.endswith("*"):
            text = part[1:-1]
            run.text = text
            run.font.color.rgb = RGBColor(*highlight_rgb)
            run.font.bold = font_bold
        else:
            run.text = part
            run.font.color.rgb = RGBColor(*default_rgb)
            run.font.bold = font_bold
        run.font.name = font_family
        run.font.size = Pt(font_size)



def populate_powerpoint_template(extracted_data, template_path, output_file, figure_paths=None, figure_descriptions=None):
    """Populate PowerPoint template with extracted academic poster information and insert up to 3 figures if provided."""
    try:
        # Load the PowerPoint template
        prs = Presentation(template_path)
        slide = prs.slides[0]
        
        # Get template name for configuration
        import os
        template_name = os.path.basename(template_path)
        template_name_without_ext = os.path.splitext(template_name)[0]
        
        # Check if this is a special template
        is_special_template = template_configs.get_template_config(template_name_without_ext) is not None
        
        # Debug output
        print(f"[DEBUG] Template name: {template_name_without_ext}")
        print(f"[DEBUG] Is special template: {is_special_template}")
        if is_special_template:
            title_settings = template_configs.get_font_settings(template_name_without_ext, "title")
            print(f"[DEBUG] Title settings: {title_settings}")
        
        # Map poster dictionary keys to shape names in template
        shape_map = {
            "headline": "HeadlineBox",
            "title": "TitleBox",
            "authors": "AuthorBox",
            "affiliations": "AffiliationBox",
            "subtitle": "SubtitleBox",
            "Introduction": "IntroductionBox",
            "Objective": "ObjectiveBox",
            "Methods": "MethodsBox",
            "Results": "ResultsBox",
            "Discussion": "DiscussionBox",
            "Conclusions": "ConclusionBox",
            "References": "ReferencesBox"
        }
        
        # Fill in the content
        title_font_size = None
        for key, shape_name in shape_map.items():
            if key == "title":
                content = extracted_data.get(key, "(No Title Extracted)")
            else:
                content = extracted_data.get(key, "")
            if content:
                # Find the shape with the exact name
                shape = find_shape_in_groups(slide, shape_name)
                
                # Add debugging for references specifically
                if key == "References":
                    print(f"🔍 [DEBUG] Looking for References shape: '{shape_name}'")
                    print(f"🔍 [DEBUG] References content: '{content[:100]}...' (length: {len(content)})")
                    if shape:
                        print(f"✅ [DEBUG] Found References shape: '{shape.name}'")
                    else:
                        print(f"❌ [DEBUG] References shape '{shape_name}' NOT found!")
                        print(f"🔍 [DEBUG] Available shapes on slide:")
                        for i, s in enumerate(slide.shapes):
                            print(f"   - Shape {i}: '{s.name}' (type: {type(s).__name__})")
                
                if shape and shape.has_text_frame:
                    if key == "headline":
                        if template_name_without_ext == "Headline Impact Template":
                            settings = template_configs.get_font_settings(template_name_without_ext, "headline")
                            font_size = settings.get("font_size", 100)
                            font_family = settings.get("font_family", "Intro Rust")
                            font_bold = settings.get("bold", True)
                            font_color = settings.get("font_color", "#ffffff")
                            alignment = settings.get("alignment", "center")
                            from pptx.util import Pt
                            from pptx.enum.text import PP_ALIGN
                            align_map = {"center": PP_ALIGN.CENTER, "left": PP_ALIGN.LEFT, "right": PP_ALIGN.RIGHT}
                            align_val = align_map.get(alignment, PP_ALIGN.CENTER)
                            highlight_rgb = (255,140,0)  # Orange
                            default_rgb = (255,255,255)  # White
                            insert_colored_headline(shape, content, highlight_rgb, default_rgb, font_family, font_size, font_bold, align_val)
                        else:
                            # fallback: treat as title
                            font_size = get_title_font_size(content, template_name_without_ext if is_special_template else None)
                            settings = template_configs.get_font_settings(template_name_without_ext, "title") if is_special_template else None
                            for paragraph in shape.text_frame.paragraphs:
                                if settings and settings.get("alignment"):
                                    alignment_constant = template_configs.get_alignment_constant(settings["alignment"])
                                    paragraph.alignment = alignment_constant
                                else:
                                    paragraph.alignment = PP_ALIGN.CENTER
                                for run in paragraph.runs:
                                    run.font.size = font_size
                                    run.font.bold = True
                                    font_family = settings.get("font_family", "Futura") if settings else "Futura"
                                    run.font.name = font_family
                                    if settings and settings.get("font_color"):
                                        rgb_color = template_configs.hex_to_rgb(settings["font_color"])
                                        run.font.color.rgb = RGBColor(*rgb_color)
                                    else:
                                        run.font.color.rgb = RGBColor(255, 255, 255)
                    else:
                        # Regular text insertion - preserves original citation format from PDF
                        shape.text = content
                        
                        # Add debugging for references insertion
                        if key == "References":
                            print(f"✅ [DEBUG] Successfully inserted references into shape '{shape.name}'")
                            print(f"✅ [DEBUG] Final references text: '{shape.text[:100]}...'")
                        # Set font size and font family based on content type and template configuration
                        if key == "title":
                            font_size = get_title_font_size(content, template_name_without_ext if is_special_template else None)
                            title_font_size = font_size.pt
                            settings = template_configs.get_font_settings(template_name_without_ext, "title") if is_special_template else None
                            for paragraph in shape.text_frame.paragraphs:
                                if settings and settings.get("alignment"):
                                    alignment_constant = template_configs.get_alignment_constant(settings["alignment"])
                                    paragraph.alignment = alignment_constant
                                else:
                                    paragraph.alignment = PP_ALIGN.CENTER
                                for run in paragraph.runs:
                                    run.font.size = font_size
                                    run.font.bold = True
                                    font_family = settings.get("font_family", "Futura") if settings else "Futura"
                                    run.font.name = font_family
                                    if settings and settings.get("font_color"):
                                        rgb_color = template_configs.hex_to_rgb(settings["font_color"])
                                        run.font.color.rgb = RGBColor(*rgb_color)
                                    else:
                                        run.font.color.rgb = RGBColor(255, 255, 255)
                        elif key == "subtitle":
                            if template_name_without_ext == "Headline Impact Template":
                                settings = template_configs.get_font_settings(template_name_without_ext, "subtitle")
                                font_size = get_subtitle_font_size(title_font_size, template_name_without_ext)
                                font_family = settings.get("font_family", "Intro Rust")
                                for paragraph in shape.text_frame.paragraphs:
                                    if settings and settings.get("alignment"):
                                        alignment_constant = template_configs.get_alignment_constant(settings["alignment"])
                                        paragraph.alignment = alignment_constant
                                    else:
                                        paragraph.alignment = PP_ALIGN.CENTER
                                    for run in paragraph.runs:
                                        run.font.size = font_size
                                        run.font.bold = False
                                        run.font.name = font_family
                                        if settings and settings.get("font_color"):
                                            rgb_color = template_configs.hex_to_rgb(settings["font_color"])
                                            run.font.color.rgb = RGBColor(*rgb_color)
                                        else:
                                            run.font.color.rgb = RGBColor(0, 0, 0)
                            else:
                                font_size = get_subtitle_font_size(title_font_size, template_name_without_ext if is_special_template else None)
                                settings = template_configs.get_font_settings(template_name_without_ext, "subtitle") if is_special_template else None
                                for paragraph in shape.text_frame.paragraphs:
                                    if settings and settings.get("alignment"):
                                        alignment_constant = template_configs.get_alignment_constant(settings["alignment"])
                                        paragraph.alignment = alignment_constant
                                    else:
                                        paragraph.alignment = PP_ALIGN.CENTER
                                    for run in paragraph.runs:
                                        run.font.size = font_size
                                        run.font.bold = False
                                        font_family = settings.get("font_family", "Futura") if settings else "Futura"
                                        run.font.name = font_family
                                        if settings and settings.get("font_color"):
                                            rgb_color = template_configs.hex_to_rgb(settings["font_color"])
                                            run.font.color.rgb = RGBColor(*rgb_color)
                                        else:
                                            run.font.color.rgb = RGBColor(0, 0, 0)
                        elif key == "authors":
                            font_size = get_dynamic_font_size(content, template_name_without_ext if is_special_template else None, section_type="authors")
                            settings = template_configs.get_font_settings(template_name_without_ext, "authors") if is_special_template else None
                            for paragraph in shape.text_frame.paragraphs:
                                if settings and settings.get("alignment"):
                                    alignment_constant = template_configs.get_alignment_constant(settings["alignment"])
                                    paragraph.alignment = alignment_constant
                                else:
                                    paragraph.alignment = PP_ALIGN.LEFT
                                for run in paragraph.runs:
                                    run.font.size = font_size
                                    font_family = settings.get("font_family", "Futura") if settings else "Futura"
                                    run.font.name = font_family
                                    run.font.bold = False
                                    if settings and settings.get("font_color"):
                                        rgb_color = template_configs.hex_to_rgb(settings["font_color"])
                                        run.font.color.rgb = RGBColor(*rgb_color)
                                    else:
                                        run.font.color.rgb = RGBColor(0, 0, 0)
                        elif key == "affiliations":
                            font_size = get_dynamic_font_size(content, template_name_without_ext if is_special_template else None, section_type="affiliations")
                            settings = template_configs.get_font_settings(template_name_without_ext, "affiliations") if is_special_template else None
                            for paragraph in shape.text_frame.paragraphs:
                                if settings and settings.get("alignment"):
                                    alignment_constant = template_configs.get_alignment_constant(settings["alignment"])
                                    paragraph.alignment = alignment_constant
                                else:
                                    paragraph.alignment = PP_ALIGN.LEFT
                                for run in paragraph.runs:
                                    run.font.size = font_size
                                    font_family = settings.get("font_family", "Futura") if settings else "Futura"
                                    run.font.name = font_family
                                    run.font.italic = settings.get("italic", False) if settings else False
                                    if settings and settings.get("font_color"):
                                        rgb_color = template_configs.hex_to_rgb(settings["font_color"])
                                        run.font.color.rgb = RGBColor(*rgb_color)
                                    else:
                                        run.font.color.rgb = RGBColor(0, 0, 0)
                        elif key == "References":
                            if template_name_without_ext == "Headline Impact Template":
                                settings = template_configs.get_font_settings(template_name_without_ext, "references")
                                font_size = Pt(settings.get("font_size", 24))
                                font_family = settings.get("font_family", "Arial")
                                for paragraph in shape.text_frame.paragraphs:
                                    if settings and settings.get("alignment"):
                                        alignment_constant = template_configs.get_alignment_constant(settings["alignment"])
                                        paragraph.alignment = alignment_constant
                                    else:
                                        paragraph.alignment = PP_ALIGN.LEFT
                                    for run in paragraph.runs:
                                        run.font.size = font_size
                                        run.font.name = font_family
                                        if settings and settings.get("font_color"):
                                            rgb_color = template_configs.hex_to_rgb(settings["font_color"])
                                            run.font.color.rgb = RGBColor(*rgb_color)
                                        else:
                                            run.font.color.rgb = RGBColor(0, 0, 0)
                            else:
                                font_size = get_dynamic_font_size(content, template_name_without_ext if is_special_template else None, section_type="references")
                                settings = template_configs.get_font_settings(template_name_without_ext, "references") if is_special_template else None
                                font_family = settings.get("font_family", "Futura") if settings else "Futura"
                                for paragraph in shape.text_frame.paragraphs:
                                    if settings and settings.get("alignment"):
                                        alignment_constant = template_configs.get_alignment_constant(settings["alignment"])
                                        paragraph.alignment = alignment_constant
                                    else:
                                        paragraph.alignment = PP_ALIGN.LEFT
                                    for run in paragraph.runs:
                                        run.font.size = font_size
                                        run.font.name = font_family
                                        if settings and settings.get("font_color"):
                                            rgb_color = template_configs.hex_to_rgb(settings["font_color"])
                                            run.font.color.rgb = RGBColor(*rgb_color)
                                        else:
                                            run.font.color.rgb = RGBColor(0, 0, 0)
                        elif key in ["summary", "executive_summary", "abstract"]:
                            font_size = Pt(24)
                            for paragraph in shape.text_frame.paragraphs:
                                for run in paragraph.runs:
                                    run.font.size = font_size
                                    run.font.name = "Futura"
                        else:
                            if template_name_without_ext == "Headline Impact Template":
                                settings = template_configs.get_font_settings(template_name_without_ext, "main_body_text")
                                font_size = Pt(settings.get("font_size", 24))
                                font_family = settings.get("font_family", "Arial")
                            else:
                                font_size = get_fixed_font_size(template_name_without_ext if is_special_template else None, section_type="main_body_text")
                                settings = template_configs.get_font_settings(template_name_without_ext, "main_body_text") if is_special_template else None
                                font_family = settings.get("font_family", "Futura") if settings else "Futura"
                            for paragraph in shape.text_frame.paragraphs:
                                if settings and settings.get("alignment"):
                                    alignment_constant = template_configs.get_alignment_constant(settings["alignment"])
                                    paragraph.alignment = alignment_constant
                                else:
                                    paragraph.alignment = PP_ALIGN.LEFT
                                for run in paragraph.runs:
                                    run.font.size = font_size
                                    run.font.name = font_family
                                    if settings and settings.get("font_color"):
                                        rgb_color = template_configs.hex_to_rgb(settings["font_color"])
                                        run.font.color.rgb = RGBColor(*rgb_color)
                                    else:
                                        run.font.color.rgb = RGBColor(0, 0, 0)
        
        # Insert up to 4 figures if provided
        if figure_paths:
            for i, fig_path in enumerate(figure_paths):
                if fig_path:
                    print(f"[DEBUG] Processing Figure {i+1}: {fig_path}")
                    
                    # Determine placeholder name based on figure number
                    # Figures 1-2: Large placeholders, Figures 3-4: Small placeholders
                    if i < 2:  # Figures 1 and 2
                        placeholder_name = f'Fig{i+1}PlaceholderLarge'
                    else:  # Figures 3, 4
                        placeholder_name = f'Fig{i+1}PlaceholderSmall'
                    
                    fig_shape = find_shape_in_groups(slide, placeholder_name)
                    if fig_shape:
                        print(f"[DEBUG] Found placeholder: {placeholder_name}")
                        
                        # Use the safe image insertion function
                        success, error = insert_image_safely(slide, fig_path, fig_shape, placeholder_name)
                        if not success:
                            print(f"[WARNING] Failed to insert image {i+1}: {error}")
                    else:
                        print(f"[DEBUG] Placeholder {placeholder_name} NOT found on slide.")
        
        # Handle figure descriptions if provided
        if figure_descriptions:
            try:
                if isinstance(figure_descriptions, str):
                    figure_descriptions = json.loads(figure_descriptions)
                
                for i in range(1, 5):  # Figures 1-4
                    if str(i) in figure_descriptions and figure_descriptions[str(i)].get('description'):
                        description = figure_descriptions[str(i)]['description']
                        if description.strip():
                            # Automatically prepend "Figure X:" to the description
                            prefixed_description = f"Figure {i}: {description.strip()}"
                            
                            # Determine description box name based on figure number
                            desc_box_name = f'FigureDesc{i}'
                            
                            desc_shape = find_shape_in_groups(slide, desc_box_name)
                            if desc_shape and desc_shape.has_text_frame:
                                # Clear existing text
                                desc_shape.text = ""
                                
                                # Apply template-specific formatting if available
                                template_name = os.path.basename(template_path)
                                template_name_without_ext = os.path.splitext(template_name)[0]
                                is_special_template = template_configs.get_template_config(template_name_without_ext) is not None
                                
                                if is_special_template:
                                    # Get font settings for figure descriptions
                                    settings = template_configs.get_font_settings(template_name_without_ext, "FigureDesc")
                                    if settings:
                                        print(f"[DEBUG] Applying FigureDesc settings for {template_name_without_ext}: {settings}")
                                        print(f"[DEBUG] Font family: {settings.get('font_family', 'Not set')}")
                                        print(f"[DEBUG] Font size: {settings.get('font_size', 'Not set')}")
                                        print(f"[DEBUG] Font color: {settings.get('font_color', 'Not set')}")
                                        
                                        # Create separate text runs for prefix (bold) and description (normal)
                                        paragraph = desc_shape.text_frame.paragraphs[0]
                                        
                                        # Set alignment
                                        if settings.get("alignment"):
                                            alignment_constant = template_configs.get_alignment_constant(settings["alignment"])
                                            paragraph.alignment = alignment_constant
                                        else:
                                            paragraph.alignment = PP_ALIGN.LEFT
                                        
                                        # Add "Figure X:" prefix in bold
                                        prefix_run = paragraph.add_run()
                                        prefix_run.text = f"Figure {i}: "
                                        font_family = settings.get("font_family", "Futura")
                                        font_size = settings.get("font_size", 18)
                                        prefix_run.font.name = font_family
                                        from pptx.util import Pt
                                        prefix_run.font.size = Pt(font_size)
                                        if settings.get("font_color"):
                                            rgb_color = template_configs.hex_to_rgb(settings["font_color"])
                                            prefix_run.font.color.rgb = RGBColor(*rgb_color)
                                        else:
                                            prefix_run.font.color.rgb = RGBColor(0, 0, 0)
                                        prefix_run.font.bold = True  # Make prefix bold
                                        
                                        print(f"[DEBUG] Figure {i} prefix: font={font_family}, size={font_size}pt, bold=True")
                                        
                                        # Add description text in normal weight
                                        desc_run = paragraph.add_run()
                                        desc_run.text = description.strip()
                                        desc_run.font.name = font_family
                                        desc_run.font.size = Pt(font_size)
                                        if settings.get("font_color"):
                                            desc_run.font.color.rgb = RGBColor(*rgb_color)
                                        else:
                                            desc_run.font.color.rgb = RGBColor(0, 0, 0)
                                        desc_run.font.bold = False  # Keep description normal
                                        
                                        print(f"[DEBUG] Figure {i} description: font={font_family}, size={font_size}pt, bold=False")
                                        
                                    else:
                                        print(f"[DEBUG] No FigureDesc settings found for {template_name_without_ext}")
                                        # Fallback: use simple text with basic formatting
                                        desc_shape.text = prefixed_description
                                else:
                                    print(f"[DEBUG] Not a special template: {template_name_without_ext}")
                                    # Fallback: use simple text with basic formatting
                                    desc_shape.text = prefixed_description
                                print(f"[DEBUG] Added description for Figure {i}: {prefixed_description[:50]}...")
                            else:
                                print(f"[DEBUG] Description box {desc_box_name} NOT found on slide.")
            except Exception as e:
                print(f"[WARNING] Error processing figure descriptions: {e}")
        

        
        # Save the presentation
        prs.save(output_file)
        return True, None
        
    except Exception as e:
        return False, f"Error populating PowerPoint template: {e}"

@app.route('/')
def index():
    """Main page with upload form."""
    return render_template('index.html')

@app.route('/test')
def test():
    """Simple test page for debugging."""
    return render_template('test.html')

@app.route('/upload', methods=['POST'])
def upload_files():
    """Handle file upload and processing."""
    try:
        print(f"[DEBUG] Upload request received. Content-Length: {request.content_length}")
        print(f"[DEBUG] Max content length: {app.config['MAX_CONTENT_LENGTH']}")
        print(f"[DEBUG] Content-Length in MB: {request.content_length // (1024*1024) if request.content_length else 0}MB")
        print(f"[DEBUG] Max content length in MB: {app.config['MAX_CONTENT_LENGTH'] // (1024*1024)}MB")
        
        # Check if request is too large before processing
        if request.content_length and request.content_length > app.config['MAX_CONTENT_LENGTH']:
            print(f"[ERROR] Request too large: {request.content_length // (1024*1024)}MB > {app.config["MAX_CONTENT_LENGTH"] // (1024*1024)}MB")
            return jsonify({
                'error': 'Total upload size too large',
                'message': f'Your upload is {request.content_length // (1024*1024)}MB, but the limit is {app.config["MAX_CONTENT_LENGTH"] // (1024*1024)}MB. Please reduce file sizes.',
                'uploaded_size_mb': request.content_length // (1024*1024),
                'max_size_mb': app.config['MAX_CONTENT_LENGTH'] // (1024*1024)
            }), 413
        
        # Get form data (academic posters only)
        extraction_type = "academic"
        global current_dummy_mode

        # Handle up to 4 figure image uploads (reduced from 6 to avoid 413 errors)
        figure_paths = [None, None, None, None]
        figures_uploaded = False
        total_figure_size = 0
        max_figure_size = 100 * 1024 * 1024  # 100MB per figure
        files_to_cleanup = []  # Track files for cleanup
        
        for i in range(1, 5):
            fig_file = request.files.get(f'figure{i}_file')
            if fig_file and fig_file.filename != '':
                if not allowed_file(fig_file.filename, {'png', 'jpg', 'jpeg'}):
                    return jsonify({'error': f'Figure {i} image must be PNG, JPG, or JPEG.'}), 400
                
                # Check individual file size
                fig_file.seek(0, 2)  # Seek to end
                file_size = fig_file.tell()
                fig_file.seek(0)  # Reset to beginning
                
                if file_size > max_figure_size:
                    return jsonify({'error': f'Figure {i} is too large ({file_size // (1024*1024)}MB). Maximum size per figure is 50MB.'}), 400
                
                total_figure_size += file_size
                if total_figure_size > MAX_CONTENT_LENGTH:
                    return jsonify({'error': f'Total figure size ({total_figure_size // (1024*1024)}MB) exceeds limit. Please reduce file sizes.'}), 400
                
                fig_path = os.path.join(UPLOAD_FOLDER, secure_filename(fig_file.filename))
                fig_file.save(fig_path)
                files_to_cleanup.append(fig_path)  # Add to cleanup list
                
                # Validate the uploaded image
                is_valid, error_msg = validate_image_file(fig_path)
                if not is_valid:
                    # Clean up the invalid file immediately
                    cleanup_uploaded_files([fig_path], keep_final_output=False)
                    return jsonify({'error': f'Figure {i}: {error_msg}'}), 400
                
                figure_paths[i-1] = fig_path
                figures_uploaded = True
        
        # Set figure_paths to None if no figures were uploaded
        if not figures_uploaded:
            figure_paths = None

        # Handle dummy mode (no PDF required)
        if current_dummy_mode:
            print("🧪 Processing in dummy mode - no PDF required")
            if not os.path.exists(DUMMY_DATA_FILE):
                return jsonify({'error': 'Dummy data not found. Please process a PDF in API mode first to create dummy data.'}), 400
            dummy_data, error = load_dummy_data()
            if error:
                return jsonify({'error': f'Error loading dummy data: {error}'}), 400
            extracted_data = dummy_data
            pdf_basename = "dummy_data"
        else:
            if 'pdf_file' not in request.files:
                return jsonify({'error': 'Please upload a PDF file for API mode.'}), 400
            pdf_file = request.files['pdf_file']
            template_file = request.files.get('template_file')
            selected_template = request.form.get('selected_template')
            if pdf_file.filename == '':
                return jsonify({'error': 'Please select a PDF file.'}), 400
            if not allowed_file(pdf_file.filename, {'pdf'}):
                return jsonify({'error': 'PDF file must have .pdf extension.'}), 400
            pdf_path = os.path.join(UPLOAD_FOLDER, secure_filename(pdf_file.filename))
            pdf_file.save(pdf_path)
            files_to_cleanup.append(pdf_path)  # Add to cleanup list
            
            manuscript_text = extract_text_from_pdf(pdf_path)
            if not manuscript_text or manuscript_text.startswith("Error"):
                cleanup_uploaded_files(files_to_cleanup, keep_final_output=False)
                return jsonify({'error': 'Failed to extract text from PDF. Please check if the PDF contains extractable text.'}), 400
            
            # Get AI provider from form data
            requested_provider = request.form.get('ai_provider', 'openai')
            print(f"🤖 User requested AI provider: {requested_provider}")
            
            # Extract information using the requested provider
            extracted_data, error = extract_information_from_pdf_with_provider(manuscript_text, requested_provider)
            if error:
                cleanup_uploaded_files(files_to_cleanup, keep_final_output=False)
                return jsonify({'error': f'Error extracting information: {error}'}), 400
            pdf_basename = os.path.splitext(os.path.basename(pdf_path))[0]

        # Template selection logic
        template_file = request.files.get('template_file')
        selected_template = request.form.get('selected_template')
        template_path = None
        if template_file and template_file.filename != '':
            if not allowed_file(template_file.filename, {'pptx'}):
                cleanup_uploaded_files(files_to_cleanup, keep_final_output=False)
                return jsonify({'error': 'Template file must have .pptx extension.'}), 400
            template_path = os.path.join(UPLOAD_FOLDER, secure_filename(template_file.filename))
            template_file.save(template_path)
            files_to_cleanup.append(template_path)  # Add to cleanup list
        elif selected_template and selected_template != 'default':
            # Search for template in all folders
            template_found = False
            for folder in ['available', 'coming_soon', 'premium']:
                folder_path = os.path.join(TEMPLATE_LIBRARY_FOLDER, folder)
                if os.path.exists(folder_path):
                    test_path = os.path.join(folder_path, selected_template)
                    if os.path.exists(test_path):
                        template_path = test_path
                        template_found = True
                        break
            
            if not template_found:
                cleanup_uploaded_files(files_to_cleanup, keep_final_output=False)
                return jsonify({'error': 'Selected template not found in library.'}), 400
        else:
            template_path = "default_template.pptx"
            if not os.path.exists(template_path):
                cleanup_uploaded_files(files_to_cleanup, keep_final_output=False)
                return jsonify({'error': 'No template selected and default template not found. Please upload a PowerPoint template or select from library.'}), 400

        # Extract figure descriptions from form data
        figure_descriptions = request.form.get('figure_descriptions', '{}')
        
        # Create output filename
        timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
        output_file = os.path.join(UPLOAD_FOLDER, f"{pdf_basename}_academic_{timestamp}.pptx")
        success, error = populate_powerpoint_template(extracted_data, template_path, output_file, figure_paths, figure_descriptions)
        if not success:
            cleanup_uploaded_files(files_to_cleanup, keep_final_output=False)
            return jsonify({'error': f'Error creating presentation: {error}'}), 400
        
        # Clean up uploaded files after successful processing
        cleanup_uploaded_files(files_to_cleanup, keep_final_output=KEEP_FINAL_OUTPUT)
        
        mode_message = "Dummy Mode" if current_dummy_mode else "API Mode"
        
        # Get AI provider information for the response (use the provider that was actually used)
        actual_provider = requested_provider if 'requested_provider' in locals() else current_api_provider
        ai_provider_info = {
            'provider': actual_provider,
            'display_name': 'ChatGPT (OpenAI)' if actual_provider == 'openai' else 'Claude (Anthropic)',
            'model': 'GPT-4o' if actual_provider == 'openai' else 'Claude 3.5 Sonnet'
        }
        
        return jsonify({
            'success': True,
            'message': f'Academic poster created successfully using {mode_message}!',
            'filename': os.path.basename(output_file),
            'extracted_data': extracted_data,
            'mode_used': mode_message,
            'ai_provider': ai_provider_info
        })
    except Exception as e:
        # Clean up any uploaded files if there was an error
        if 'files_to_cleanup' in locals():
            cleanup_uploaded_files(files_to_cleanup, keep_final_output=False)
        return jsonify({'error': f'An unexpected error occurred: {e}'}), 500

@app.route('/download/<filename>')
def download_file(filename):
    """Download the generated PowerPoint file and optionally clean it up."""
    try:
        file_path = os.path.join(UPLOAD_FOLDER, filename)
        if os.path.exists(file_path):
            # If auto-cleanup is enabled, delete the file after download
            if AUTO_CLEANUP_UPLOADS and KEEP_FINAL_OUTPUT:
                # Use a callback to delete the file after sending
                def cleanup_after_send():
                    try:
                        if os.path.exists(file_path):
                            os.remove(file_path)
                            print(f"🗑️ Cleaned up after download: {filename}")
                    except Exception as e:
                        print(f"⚠️ Could not delete file after download {filename}: {e}")
                
                response = send_file(file_path, as_attachment=True, download_name=filename)
                response.call_on_close(cleanup_after_send)
                return response
            else:
                return send_file(file_path, as_attachment=True, download_name=filename)
        else:
            return jsonify({'error': 'File not found.'}), 404
    except Exception as e:
        return jsonify({'error': f'Error downloading file: {e}'}), 500

@app.route('/api/mode', methods=['GET'])
def get_mode():
    """Get current dummy data mode."""
    global current_dummy_mode
    return jsonify({
        'use_dummy_data': current_dummy_mode,
        'mode_name': 'Dummy Mode' if current_dummy_mode else 'API Mode'
    })

@app.route('/api/mode', methods=['POST'])
def set_mode():
    """Set dummy data mode."""
    global current_dummy_mode
    try:
        data = request.get_json()
        if data is None:
            return jsonify({'success': False, 'error': 'No JSON data provided'}), 400
        
        use_dummy_data = data.get('use_dummy_data', False)
        current_dummy_mode = use_dummy_data
        
        print(f"🔄 Mode changed to: {'Dummy Mode' if use_dummy_data else 'API Mode'}")
        
        return jsonify({
            'success': True,
            'use_dummy_data': current_dummy_mode,
            'message': f"Switched to {'Dummy Mode' if use_dummy_data else 'API Mode'}"
        })
        
    except Exception as e:
        return jsonify({'success': False, 'error': str(e)}), 500

@app.route('/api/provider', methods=['GET'])
def get_api_provider():
    """Get current API provider."""
    global current_api_provider
    return jsonify({
        'provider': current_api_provider,
        'available_providers': {
            'openai': {
                'name': 'ChatGPT (OpenAI)',
                'configured': bool(OPENAI_API_KEY),
                'model': 'GPT-4o'
            },
            'anthropic': {
                'name': 'Claude (Anthropic)',
                'configured': bool(ANTHROPIC_API_KEY),
                'model': 'Claude 3.5 Sonnet'
            }
        }
    })

@app.route('/api/provider', methods=['POST'])
def set_api_provider():
    """Set API provider."""
    global current_api_provider
    try:
        data = request.get_json()
        if data is None:
            return jsonify({'success': False, 'error': 'No JSON data provided'}), 400
        
        provider = data.get('provider', 'openai')
        
        # Validate provider
        if provider not in ['openai', 'anthropic']:
            return jsonify({'success': False, 'error': 'Invalid provider. Must be "openai" or "anthropic"'}), 400
        
        # Check if API key is configured
        if provider == 'openai' and not OPENAI_API_KEY:
            return jsonify({'success': False, 'error': 'OpenAI API key not configured'}), 400
        elif provider == 'anthropic' and not ANTHROPIC_API_KEY:
            return jsonify({'success': False, 'error': 'Anthropic API key not configured'}), 400
        
        current_api_provider = provider
        
        provider_names = {
            'openai': 'ChatGPT (OpenAI)',
            'anthropic': 'Claude (Anthropic)'
        }
        
        print(f"🔄 API Provider changed to: {provider_names[provider]}")
        
        return jsonify({
            'success': True,
            'provider': current_api_provider,
            'message': f"Switched to {provider_names[provider]}"
        })
        
    except Exception as e:
        return jsonify({'success': False, 'error': str(e)}), 500

@app.route('/api/dummy-data-status')
def get_dummy_data_status():
    """Get dummy data availability and preview."""
    try:
        if os.path.exists(DUMMY_DATA_FILE):
            with open(DUMMY_DATA_FILE, 'r', encoding='utf-8') as f:
                dummy_data = json.load(f)
            
            # Get a preview of the data (first 100 chars of each field)
            preview = {}
            for key, value in dummy_data['data'].items():
                if isinstance(value, str):
                    preview[key] = value[:100] + ('...' if len(value) > 100 else '')
                else:
                    preview[key] = str(value)
            
            return jsonify({
                'available': True,
                'timestamp': dummy_data['timestamp'],
                'preview': preview,
                'fields': list(dummy_data['data'].keys())
            })
        else:
            return jsonify({
                'available': False,
                'message': 'No dummy data file found'
            })
            
    except Exception as e:
        return jsonify({
            'available': False,
            'error': str(e)
        })

@app.route('/api/template-library')
def get_template_library():
    """Get list of templates in the library."""
    try:
        print(f"[DEBUG] Template library API called")
        print(f"[DEBUG] Template library folder: {TEMPLATE_LIBRARY_FOLDER}")
        print(f"[DEBUG] Folder exists: {os.path.exists(TEMPLATE_LIBRARY_FOLDER)}")
        
        templates = load_template_library()
        print(f"[DEBUG] load_template_library() returned {len(templates)} templates")
        
        return jsonify({
            'success': True,
            'templates': templates,
            'count': len(templates)
        })
    except Exception as e:
        print(f"[DEBUG] Error in template library API: {e}")
        import traceback
        traceback.print_exc()
        return jsonify({
            'success': False,
            'error': str(e)
        }), 500

@app.route('/api/template-library/upload', methods=['POST'])
def upload_template_to_library():
    """Upload a template to the library."""
    try:
        if 'template_file' not in request.files:
            return jsonify({'success': False, 'error': 'No template file provided'}), 400
        
        template_file = request.files['template_file']
        preview_file = request.files.get('preview_file')
        if template_file.filename == '':
            return jsonify({'success': False, 'error': 'No template file selected'}), 400
        
        if not allowed_file(template_file.filename, {'pptx'}):
            return jsonify({'success': False, 'error': 'Template file must have .pptx extension'}), 400
        
        # Save to library
        filename = secure_filename(template_file.filename)
        success, error = save_template_to_library(template_file, filename)
        
        # Handle preview image if provided
        preview_filename = None
        if preview_file and preview_file.filename != '':
            ext = os.path.splitext(preview_file.filename)[1].lower()
            if ext not in ['.png', '.jpg', '.jpeg']:
                return jsonify({'success': False, 'error': 'Preview image must be PNG or JPG'}), 400
            preview_filename = os.path.splitext(filename)[0] + '_manual_preview' + ext
            preview_path = os.path.join(TEMPLATE_LIBRARY_FOLDER, preview_filename)
            preview_file.save(preview_path)
        
        if success:
            return jsonify({
                'success': True,
                'message': f'Template "{filename}" added to library successfully!',
                'filename': filename
            })
        else:
            return jsonify({'success': False, 'error': error}), 500
            
    except Exception as e:
        return jsonify({'success': False, 'error': str(e)}), 500

@app.route('/api/template-library/<filename>', methods=['DELETE'])
def delete_template_from_library(filename):
    """Delete a template from the library."""
    try:
        template_path = os.path.join(TEMPLATE_LIBRARY_FOLDER, filename)
        if os.path.exists(template_path):
            os.remove(template_path)
            
            # Also delete preview if it exists
            preview_filename = os.path.splitext(filename)[0] + '_preview.png'
            preview_path = os.path.join(TEMPLATE_LIBRARY_FOLDER, preview_filename)
            if os.path.exists(preview_path):
                os.remove(preview_path)
            
            # Also delete manual preview if it exists
            base = os.path.splitext(filename)[0]
            for ext in ['.png', '.jpg', '.jpeg']:
                manual_preview_path = os.path.join(TEMPLATE_LIBRARY_FOLDER, base + '_manual_preview' + ext)
                if os.path.exists(manual_preview_path):
                    os.remove(manual_preview_path)
            
            return jsonify({
                'success': True,
                'message': f'Template "{filename}" deleted from library'
            })
        else:
            return jsonify({'success': False, 'error': 'Template not found'}), 404
            
    except Exception as e:
        return jsonify({'success': False, 'error': str(e)}), 500

@app.route('/api/template-library/preview/<filename>')
def get_template_preview(filename):
    """Get template preview image."""
    try:
        # Check for manual preview first in all folders
        for folder in ['available', 'coming_soon', 'premium']:
            folder_path = os.path.join(TEMPLATE_LIBRARY_FOLDER, folder)
            if os.path.exists(folder_path):
                # Check if the exact filename exists
                exact_preview_path = os.path.join(folder_path, filename)
                if os.path.exists(exact_preview_path):
                    return send_file(exact_preview_path, mimetype='image/png')
                
                # Check for manual preview variants
                base = os.path.splitext(filename)[0]
                for ext in ['.png', '.jpg', '.jpeg']:
                    manual_preview_path = os.path.join(folder_path, base + '_manual_preview' + ext)
                    if os.path.exists(manual_preview_path):
                        return send_file(manual_preview_path, mimetype='image/png')
        
        return jsonify({'error': 'Preview not found'}), 404
    except Exception as e:
        return jsonify({'error': f'Error serving preview: {e}'}), 500

@app.route('/api/template-library/delete', methods=['POST'])
def delete_template():
    """Move a template to the archive instead of deleting it."""
    try:
        data = request.get_json()
        filename = data.get('filename')
        
        if not filename:
            return jsonify({'success': False, 'error': 'No filename provided'}), 400
        
        # Prevent deletion of default template
        if filename == 'default':
            return jsonify({'success': False, 'error': 'Cannot delete default template'}), 400
        
        # Create archive folder if it doesn't exist
        archive_folder = os.path.join(os.path.dirname(TEMPLATE_LIBRARY_FOLDER), 'Template Archive')
        if not os.path.exists(archive_folder):
            os.makedirs(archive_folder)
            print(f"📁 Created archive folder: {archive_folder}")
        
        # Search for template in all folders
        template_found = False
        template_folder = None
        template_path = None
        
        for folder in ['available', 'coming_soon', 'premium']:
            folder_path = os.path.join(TEMPLATE_LIBRARY_FOLDER, folder)
            if os.path.exists(folder_path):
                test_path = os.path.join(folder_path, filename)
                if os.path.exists(test_path):
                    template_path = test_path
                    template_folder = folder
                    template_found = True
                    break
        
        if not template_found:
            return jsonify({'success': False, 'error': 'Template not found'}), 404
        
        base = os.path.splitext(filename)[0]
        
        # Move the template file to archive
        archive_template_path = os.path.join(archive_folder, filename)
        import shutil
        shutil.move(template_path, archive_template_path)
        print(f"📦 Archived template: {filename} from {template_folder}")
        
        # Move associated preview files to archive
        preview_extensions = ['.png', '.jpg', '.jpeg']
        moved_files = [filename]  # Track all moved files
        
        for ext in preview_extensions:
            # Move auto-generated preview
            preview_path = os.path.join(TEMPLATE_LIBRARY_FOLDER, template_folder, base + '_preview' + ext)
            if os.path.exists(preview_path):
                archive_preview_path = os.path.join(archive_folder, base + '_preview' + ext)
                shutil.move(preview_path, archive_preview_path)
                moved_files.append(base + '_preview' + ext)
                print(f"📦 Archived preview: {base}_preview{ext}")
            
            # Move manual preview
            manual_preview_path = os.path.join(TEMPLATE_LIBRARY_FOLDER, template_folder, base + '_manual_preview' + ext)
            if os.path.exists(manual_preview_path):
                archive_manual_preview_path = os.path.join(archive_folder, base + '_manual_preview' + ext)
                shutil.move(manual_preview_path, archive_manual_preview_path)
                moved_files.append(base + '_manual_preview' + ext)
                print(f"📦 Archived manual preview: {base}_manual_preview{ext}")
        
        return jsonify({
            'success': True, 
            'message': f'Template {filename} moved to archive successfully',
            'archived_files': moved_files
        })
        
    except Exception as e:
        print(f"❌ Error archiving template: {e}")
        return jsonify({'success': False, 'error': str(e)}), 500

@app.route('/api/premium-templates', methods=['GET'])
def get_premium_templates():
    """Get list of premium templates."""
    try:
        from template_configs import get_premium_templates
        premium_templates = get_premium_templates()
        return jsonify({
            'success': True,
            'premium_templates': premium_templates,
            'count': len(premium_templates)
        })
    except Exception as e:
        return jsonify({
            'success': False,
            'error': str(e)
        }), 500

@app.route('/api/premium-templates', methods=['POST'])
def add_premium_template():
    """Add a template to premium list."""
    try:
        from template_configs import add_premium_template
        data = request.get_json()
        if not data or 'template_name' not in data:
            return jsonify({'success': False, 'error': 'Template name required'}), 400
        
        template_name = data['template_name']
        success = add_premium_template(template_name)
        
        if success:
            return jsonify({
                'success': True,
                'message': f'Template "{template_name}" added to premium list'
            })
        else:
            return jsonify({
                'success': False,
                'error': f'Template "{template_name}" is already premium'
            }), 400
            
    except Exception as e:
        return jsonify({'success': False, 'error': str(e)}), 500

@app.route('/api/premium-templates', methods=['DELETE'])
def remove_premium_template():
    """Remove a template from premium list."""
    try:
        from template_configs import remove_premium_template
        data = request.get_json()
        if not data or 'template_name' not in data:
            return jsonify({'success': False, 'error': 'Template name required'}), 400
        
        template_name = data['template_name']
        success = remove_premium_template(template_name)
        
        if success:
            return jsonify({
                'success': True,
                'message': f'Template "{template_name}" removed from premium list'
            })
        else:
            return jsonify({
                'success': False,
                'error': f'Template "{template_name}" is not in premium list'
            }), 400
            
    except Exception as e:
        return jsonify({'success': False, 'error': str(e)}), 500

@app.route('/api/coming-soon-templates', methods=['GET'])
def get_coming_soon_templates():
    """Get list of coming soon templates."""
    try:
        from template_configs import get_coming_soon_templates
        coming_soon_templates = get_coming_soon_templates()
        return jsonify({
            'success': True,
            'coming_soon_templates': coming_soon_templates,
            'count': len(coming_soon_templates)
        })
    except Exception as e:
        return jsonify({
            'success': False,
            'error': str(e)
        }), 500

@app.route('/api/coming-soon-templates', methods=['POST'])
def add_coming_soon_template():
    """Add a template to coming soon list."""
    try:
        from template_configs import add_coming_soon_template
        data = request.get_json()
        if not data or 'template_name' not in data:
            return jsonify({'success': False, 'error': 'Template name required'}), 400
        
        template_name = data['template_name']
        success = add_coming_soon_template(template_name)
        
        if success:
            return jsonify({
                'success': True,
                'message': f'Template "{template_name}" added to coming soon list'
            })
        else:
            return jsonify({
                'success': False,
                'error': f'Template "{template_name}" is already in coming soon list'
            }), 400
            
    except Exception as e:
        return jsonify({'success': False, 'error': str(e)}), 500

@app.route('/api/coming-soon-templates', methods=['DELETE'])
def remove_coming_soon_template():
    """Remove a template from coming soon list."""
    try:
        from template_configs import remove_coming_soon_template
        data = request.get_json()
        if not data or 'template_name' not in data:
            return jsonify({'success': False, 'error': 'Template name required'}), 400
        
        template_name = data['template_name']
        success = remove_coming_soon_template(template_name)
        
        if success:
            return jsonify({
                'success': True,
                'message': f'Template "{template_name}" removed from coming soon list'
            })
        else:
            return jsonify({
                'success': False,
                'error': f'Template "{template_name}" is not in coming soon list'
            }), 400
            
    except Exception as e:
        return jsonify({'success': False, 'error': str(e)}), 500

@app.route('/api/template-descriptions')
def get_template_descriptions():
    """Get all template descriptions."""
    try:
        from template_configs import get_all_template_descriptions
        descriptions = get_all_template_descriptions()
        return jsonify({
            'success': True,
            'descriptions': descriptions
        })
    except Exception as e:
        return jsonify({'success': False, 'error': str(e)}), 500

@app.route('/api/template-descriptions/<template_name>')
def get_template_description(template_name):
    """Get description for a specific template."""
    try:
        from template_configs import get_template_description
        description = get_template_description(template_name)
        
        if description:
            return jsonify({
                'success': True,
                'description': description
            })
        else:
            return jsonify({
                'success': False,
                'error': 'Template description not found'
            }), 404
            
    except Exception as e:
        return jsonify({'success': False, 'error': str(e)}), 500

@app.errorhandler(413)
def too_large(e):
    """Handle file too large error."""
    return jsonify({
        'error': 'File too large',
        'message': 'The total size of uploaded files exceeds the limit. Please reduce file sizes or upload fewer files.',
        'max_size_mb': MAX_CONTENT_LENGTH // (1024 * 1024),
        'individual_figure_limit_mb': 100,
        'max_figures': 4,
        'tips': [
            'Compress your images before uploading',
            'Use JPG format instead of PNG for photos',
            'Reduce image resolution if possible',
            'Upload fewer figures if needed'
        ]
    }), 413

@app.route('/health')
def health_check():
    """Health check endpoint."""
    return jsonify({'status': 'healthy', 'timestamp': datetime.now().isoformat()})

@app.route('/api/upload-limits')
def get_upload_limits():
    """Get current upload limits for debugging."""
    return jsonify({
        'max_content_length_mb': MAX_CONTENT_LENGTH // (1024 * 1024),
        'max_figure_size_mb': 100,
        'max_figures': 4,
        'allowed_extensions': list(ALLOWED_EXTENSIONS),
        'tips': [
            'If you get 413 errors, try compressing your images',
            'Use JPG instead of PNG for photos',
            'Reduce image resolution to 300 DPI or less',
            'Upload fewer figures if needed'
        ]
    })

@app.route('/api/test-upload', methods=['POST'])
def test_upload():
    """Test endpoint to check if uploads work at all."""
    try:
        print(f"[DEBUG] Test upload received. Content-Length: {request.content_length}")
        print(f"[DEBUG] Max content length: {app.config['MAX_CONTENT_LENGTH']}")
        print(f"[DEBUG] Content-Length in MB: {request.content_length // (1024*1024) if request.content_length else 0}MB")
        
        # Check if request is too large
        if request.content_length and request.content_length > app.config['MAX_CONTENT_LENGTH']:
            return jsonify({
                'error': 'Request too large',
                'uploaded_size_mb': request.content_length // (1024*1024),
                'max_size_mb': app.config['MAX_CONTENT_LENGTH'] // (1024*1024)
            }), 413
        
        # Check files
        files_info = []
        total_size = 0
        
        for key, file in request.files.items():
            if file and file.filename:
                file.seek(0, 2)  # Seek to end
                file_size = file.tell()
                file.seek(0)  # Reset to beginning
                total_size += file_size
                
                files_info.append({
                    'name': file.filename,
                    'size_mb': file_size // (1024*1024),
                    'type': file.content_type
                })
        
        return jsonify({
            'success': True,
            'message': 'Test upload successful',
            'total_size_mb': total_size // (1024*1024),
            'files': files_info,
            'max_size_mb': app.config['MAX_CONTENT_LENGTH'] // (1024*1024)
        })
        
    except Exception as e:
        return jsonify({
            'error': f'Test upload failed: {str(e)}',
            'exception_type': type(e).__name__
        }), 500

@app.route('/api/cleanup', methods=['POST'])
def manual_cleanup():
    """Manually trigger cleanup of old files."""
    try:
        data = request.get_json() or {}
        days_old = data.get('days_old', 1)
        
        # Count files before cleanup
        files_before = len([f for f in os.listdir(UPLOAD_FOLDER) if os.path.isfile(os.path.join(UPLOAD_FOLDER, f))])
        
        # Run cleanup
        cleanup_old_files(days_old)
        
        # Count files after cleanup
        files_after = len([f for f in os.listdir(UPLOAD_FOLDER) if os.path.isfile(os.path.join(UPLOAD_FOLDER, f))])
        files_removed = files_before - files_after
        
        return jsonify({
            'success': True,
            'message': f'Cleanup completed. Removed {files_removed} old files.',
            'files_before': files_before,
            'files_after': files_after,
            'files_removed': files_removed
        })
        
    except Exception as e:
        return jsonify({'error': f'Cleanup failed: {e}'}), 500

@app.route('/api/cleanup-status')
def cleanup_status():
    """Get current cleanup settings and upload folder status."""
    try:
        files_in_uploads = [f for f in os.listdir(UPLOAD_FOLDER) if os.path.isfile(os.path.join(UPLOAD_FOLDER, f))]
        total_size = sum(os.path.getsize(os.path.join(UPLOAD_FOLDER, f)) for f in files_in_uploads)
        
        return jsonify({
            'auto_cleanup_enabled': AUTO_CLEANUP_UPLOADS,
            'keep_final_output': KEEP_FINAL_OUTPUT,
            'files_count': len(files_in_uploads),
            'total_size_mb': round(total_size / (1024 * 1024), 2),
            'files': files_in_uploads[:10]  # Show first 10 files
        })
        
    except Exception as e:
        return jsonify({'error': f'Could not get cleanup status: {e}'}), 500

if __name__ == '__main__':
    print("🚀 Starting Academic Poster Creator...")
    
    # Get port from environment variable (for deployment) or use default
    port = int(os.environ.get('PORT', 5000))
    
    print(f"📱 Open your browser and go to: http://localhost:{port}")
    print(f"📏 Total upload limit: {MAX_CONTENT_LENGTH // (1024*1024)}MB")
    print(f"📏 Individual figure limit: 100MB per figure")
    print(f"📏 Maximum figures: 4 figures (reduced to avoid 413 errors)")
    
    # Show API configuration status
    print("\n🤖 AI API Configuration:")
    if OPENAI_API_KEY:
        print("✅ OpenAI API key configured")
    else:
        print("❌ OpenAI API key not configured")
    
    if ANTHROPIC_API_KEY:
        print("✅ Anthropic API key configured")
    else:
        print("❌ Anthropic API key not configured")
    
    print(f"🔧 Default API provider: {current_api_provider.upper()}")
    print(f"🔧 Anthropic version: {anthropic.__version__}")
    print(f"🔧 Environment check - OpenAI key: {'✅ Set' if OPENAI_API_KEY else '❌ Missing'}")
    if OPENAI_API_KEY:
        print(f"🔧 OpenAI key starts with: {OPENAI_API_KEY[:10]}...")
    print(f"🔧 Environment check - Anthropic key: {'✅ Set' if ANTHROPIC_API_KEY else '❌ Missing'}")
    if ANTHROPIC_API_KEY:
        print(f"🔧 Anthropic key starts with: {ANTHROPIC_API_KEY[:10]}...")
    
    # Check for proxy-related environment variables that might cause issues
    proxy_vars = ['HTTP_PROXY', 'HTTPS_PROXY', 'http_proxy', 'https_proxy']
    for var in proxy_vars:
        if os.getenv(var):
            print(f"⚠️ Found proxy environment variable: {var}={os.getenv(var)}")
        else:
            print(f"✅ No proxy variable: {var}")
    
    # Clean up old files on startup if auto-cleanup is enabled
    if AUTO_CLEANUP_UPLOADS:
        print("🧹 Auto-cleanup enabled - cleaning old files on startup...")
        cleanup_old_files(1)  # Clean files older than 1 day
        print("✅ Startup cleanup completed")
    
    if current_dummy_mode:
        print("🧪 Using dummy data for testing (no API calls)")
    else:
        provider_names = {
            'openai': 'OpenAI (ChatGPT)',
            'anthropic': 'Anthropic (Claude)'
        }
        print(f"💡 Using {provider_names.get(current_api_provider, current_api_provider.upper())} API for intelligent information extraction!")
        print("🔑 Make sure your API keys are set in .env file")
    
    # Use environment variable for debug mode (False in production)
    debug_mode = os.environ.get('FLASK_DEBUG', 'False').lower() == 'true'
    app.run(debug=debug_mode, host='0.0.0.0', port=port) 